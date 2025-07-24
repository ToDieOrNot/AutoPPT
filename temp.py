#!/usr/bin/env python
# -*- coding: UTF-8 -*-
"""
@Project ：AutoPPT
@File    ：page_prompt.py
@Date    ：2025/7/22
@Descrip ：
"""
from views.page_public import *
# from models.view_prompts_curd import *
# from models.view_template_pptx_curd import *
# from views.page_template_pptx import *
from models.config_read import read_env, read_json_file

import pywebio
from pywebio import config
from pywebio.input import *
from pywebio.output import *
from pywebio.session import *
from pywebio.platform.tornado import start_server
from models.config_read import read_env
from models.config_read import recursive_list_files, read_json_file
from models import view_models_curd
from models import view_run
from views import page_run

model_res = """

{
    "1": {
        "name": "成果",
        "type": "内容页",
        "descript": "展示团队目前已取得的成果",
        "first_title1": "1、模型试错流程",
        "first_title2": "",
        "second_title1": "2、四川方言识别方案",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "目前，我们团队进行了大量模型试错，完成了设计、开发、部署、训练、推理的流程。",
        "context2": "更新了四川方言的识别方案。",
        "pic1": "",
        "pic2": ""
    },
    "2": {
        "name": "成果",
        "type": "内容页",
        "descript": "继续展示团队成果",
        "first_title1": "3、AI外呼经验",
        "first_title2": "",
        "second_title1": "4、辅助谭总",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "学习了'反诈中心'的AI外呼经验，结合我们实际，研究出'语音知识库'的高效方案。",
        "context2": "（1）提供了基本问答库、（2）提供了训练模型数据、（3）外呼通知、（4）微信机器人。",
        "pic1": "",
        "pic2": ""
    },
    "3": {
        "name": "成果",
        "type": "内容页",
        "descript": "继续展示团队成果",
        "first_title1": "5、蓉智考系统",
        "first_title2": "",
        "second_title1": "6、协助工作",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "（1）目前AI的使用仍停留在初级阶段，即调用豆包的服务进行问答和部署基础模型。（2）外呼。使用开发板连接手机卡，完成'拨号-发送录音-挂断'的流程。",
        "context2": "（1）协议送簽。（2）部分录课。（3）其他协助。",
        "pic1": "",
        "pic2": ""
    },
    "4": {
        "name": "存在问题",
        "type": "内容页",
        "descript": "列出当前存在的问题",
        "first_title1": "1、手机和服务端通信信号的双向传输",
        "first_title2": "",
        "second_title1": "2、主动引导型话术资源的建设",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "即实现主机的听和说。",
        "context2": "叶老师一直比较忙，沟通较少。向孔老师学习话术时提到话术体系的复杂性，如催费前要在前几次嘘寒问暖中铺垫；替考需要从学生实际情况考虑。",
        "pic1": "",
        "pic2": ""
    },
    "5": {
        "name": "存在问题",
        "type": "内容页",
        "descript": "继续列出当前存在的问题",
        "first_title1": "3、研究和测试需要基础设备",
        "first_title2": "",
        "second_title1": "",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "我们部署的Mac mini M4(3000元左右)不支持大部分的语音大模型。",
        "context2": "",
        "pic1": "",
        "pic2": ""
    },
    "6": {
        "name": "解决方案",
        "type": "内容页",
        "descript": "介绍可能的解决方案",
        "first_title1": "1、通话信号双向传输",
        "first_title2": "",
        "second_title1": "2、话术体系建设",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "（1）购买云联络平台，使用平台虚拟号。（示例：Twilio）（2）购买VTU+SIP方案，使用自己的手机卡。（流程：SIP服务器收到呼叫指令后会自动控制VTU设备呼叫指定手机号，并可将对方是否接通的信号通过websocket或MRCP协议反馈。）（3）专业的AI外呼开发人员。（4）备选方案：使用电话机控制拨号和挂断；配合电脑的外设，音箱和话筒；需要隔音箱，同时保留散热通道。",
        "context2": "（1）监听公用手机设备，实时显示话术提示和参考资料，收集对话数据，统计KPI。（2）专人对历史对话数据进行声纹识别，按场景、对话人、实际效果进行评定。",
        "pic1": "",
        "pic2": ""
    },
    "7": {
        "name": "解决方案",
        "type": "内容页",
        "descript": "继续介绍解决方案",
        "first_title1": "3、设备",
        "first_title2": "",
        "second_title1": "3.1 基础设备-V100(开发环境)",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "配置：使用E5处理器 & X99主板 & 32G内存 & 24GV100 & 1T硬盘。优点：能运行大部分语音模型。缺点：较大参数的模型需要降低模型精度使用，不能同时启用多个模型（即'听-思考-说'的通话延迟很高）。成本：价格范围在3500元-8000元，一月耗电720度左右。完整的服务器较少，需要自己组装。",
        "context2": "3.2 租赁设备（开发环境）。优点：成本较低，可随时关闭。缺点：（1）端口有限。（2）不能与本地手机端进行通信。（3）虚拟显卡存在显著的延迟性。成本：每月2000元左右。3.3 实际生产设备（生产环境）。显卡：需要支持cuda12.8 并且 FP32应不低于40TFLOPS 并且 显存不低于32G。内存：应大于等于显存。处理器：不少于12核。优点：可支持1-2个实时对话，并且拥有较低的延迟性。成本：价格范围在20000元以内。", 
        "pic1": "",
        "pic2": ""
    },
    "8": {
        "name": "待完成工作内容",
        "type": "内容页",
        "descript": "展示待完成的工作内容",
        "first_title1": "1、蓉智考上线后测试",
        "first_title2": "",
        "second_title1": "2、学生界面原型重构",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "（不含负载测试、单元测试等）",
        "context2": "",
        "pic1": "",
        "pic2": ""
    },
    "9": {
        "name": "待完成工作内容",
        "type": "内容页",
        "descript": "继续介绍待完成的工作内容",
        "first_title1": "3、统计各部门使用问题及协商优化",
        "first_title2": "",
        "second_title1": "4、教务话术体系在蓉智考平台的使用",
        "second_title2": "",
        "third_title1": "",
        "third_title2": "",
        "context1": "确认签字",
        "context2": "真题自动录入，教育政策变化监听",
        "pic1": "",
        "pic2": ""
    }
}
"""

model_res1 = {
  "1": {
    "name": "首页",
    "type": "章节主页",
    "descript": "项目总结与成果展示",
    "first_title1": "项目成果",
    "first_title2": "",
    "second_title1": "",
    "second_title2": "",
    "third_title1": "",
    "third_title2": "",
    "context1": "我们团队完成了模型试错、开发、部署、训练和推理的全流程，同时优化了四川方言的识别方案，并借鉴了反诈中心的AI外呼经验，得到了高效的语音知识库方案。",
    "context2": "",
    "pic1": "",
    "pic2": ""
  },
  "2": {
    "name": "项目成果",
    "type": "小标题正文",
    "descript": "项目当前的完成情况",
    "first_title1": "项目成果",
    "first_title2": "",
    "second_title1": "1. 完成全流程",
    "second_title2": "2. 方言识别优化",
    "third_title1": "3. 语音知识库方案",
    "third_title2": "4. 协助谭总工作",
    "context1": "目前AI的使用仍停留在初级阶段，调用豆包的服务进行问答和部署基础模型。通话流程方面，完成了拨号-发送录音-挂断的流程。",
    "context2": "蓉智考系统当前还在使用中，学生界面原型也在进行重构。",
    "pic1": "",
    "pic2": ""
  },
  "3": {
    "name": "存在问题",
    "type": "标题和内容",
    "descript": "目前存在的主要问题",
    "first_title1": "问题1: 通信双向传输",
    "first_title2": "问题2: 话术资源建设",
    "second_title1": "问题3: 设备限制",
    "second_title2": "问题4: 未完成的工作",
    "third_title1": "",
    "third_title2": "",
    "context1": "手机和服务端之间的双向通信尚未实现，限制了主持人的听和说功能。此外，话术资源的建设需要更多的资源和时间，叶老师工作繁忙，沟通较少。",
    "context2": "目前基础设备不支持大部分语音大模型，且需要自己组装。未完成的工作包括蓉智考的测试、学生功能界面的重构、部门问题协商及优化等。",
    "pic1": "",
    "pic2": ""
  },
  "4": {
    "name": "解决方案",
    "type": "标题和内容",
    "descript": "针对当前问题的解决方案",
    "first_title1": "解决方案1: 双向通信",
    "first_title2": "解决方案2: 话术体系",
    "second_title1": "解决方案3: 设备升级",
    "second_title2": "解决方案4: 线上服务",
    "third_title1": "",
    "third_title2": "",
    "context1": "考虑购买云联络平台或VTU+SIP方案支持双向通信，目前也在考虑外呼的备选方案，如电话机配合电脑外设。",
    "context2": "通过专人分析历史对话数据进行声纹识别并按场景、对话人、实际效果进行评定，同时利用电脑实时显示和收集数据。",
    "pic1": "",
    "pic2": ""
  },
  "5": {
    "name": "设备升级",
    "type": "内容页",
    "descript": "设备升级与选择",
    "first_title1": "基础设备-V100",
    "first_title2": "租赁设备",
    "second_title1": "实际生产设备",
    "second_title2": "购买线上服务",
    "third_title1": "",
    "third_title2": "",
    "context1": "基础设备V100配备了E5处理器、X99主板、32G内存和24GV100显卡，可运行大部分语音模型，但无法同时启用多个模型，延迟较高。",
    "context2": "租赁设备有成本低、易关闭的优点，但存在端口限制、无法与本地通信和延迟高的问题。",
    "pic1": "",
    "pic2": ""
  },
  "6": {
    "name": "待完成工作",
    "type": "内容页",
    "descript": "待完成的主要任务",
    "first_title1": "蓉智考上线测试",
    "first_title2": "学生界面重构",
    "second_title1": "问题协商和优化",
    "second_title2": "话术体系应用",
    "third_title1": "教研内容",
    "third_title2": "",
    "context1": "蓉智考系统上线后的测试需要完成，包括接口功能、用户体验等。同时需重构学生界面原型，组织多部门解决问题并确认签字。",
    "context2": "教务话术体系需要在蓉智考平台上进行应用。此外，需要加入真题自动录入以及教育政策变化监听的功能。",
    "pic1": "",
    "pic2": ""
  }
}


pptx_demo = "./files/templates/pptx\demo.pptx"

def start():
    page_run.load_interface()


import os, json, uuid
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

def output_pptx(aimodel_res, pptx_demo):
    dict_env = read_env()
    dir_pptx = dict_env.get("file_json_pptx")
    dir_output = dict_env.get("dir_output_pptx")
    filename = os.path.splitext(os.path.basename(os.path.normpath(pptx_demo)))[0]
    output_filepath = os.path.join(dir_output, filename + "_" + str(datetime.now().strftime("%Y%m%d%H%M%S")) + ".pptx")
    pptx_json = read_json_file(dir_pptx).get(filename)
    format_pptx_json = {j["type"]: {**j, "page": i} for i, j in pptx_json.items()}
    valid_types = set(format_pptx_json.keys())  # 获取所有有效的PPT类型

    prs_source = Presentation(pptx_demo)
    prs_target = Presentation()
    for ri, rj in aimodel_res.items():
        if rj["type"] not in valid_types:
            continue  # 如果当前页的类型不在有效类型列表中，则跳过该页
        page_pptx_type = format_pptx_json[rj["type"]]
        reversed_page_pptx_type = {pj: pi for pi, pj in page_pptx_type.items() if pj and (("title" in pi) or ("context" in pi) or ("pic" in pi))}
        page_num = int(page_pptx_type["page"]) - 1

        template_slide = prs_source.slides[page_num]
        new_slide = prs_target.slides.add_slide(template_slide)
        # 复制所有形状（文本、图像等）
        for shp in prs_source.slides[page_num].shapes:
            if shp.has_text_frame:
                # 复制文本框内容
                idx = len(new_slide.shapes)
                el = shp.element
                new_slide.shapes._spTree.insert(idx, el)
            elif shp.has_table:
                # 复制表格（保持结构和内容）
                idx = len(new_slide.shapes)
                el = shp.element
                new_slide.shapes._spTree.insert(idx, el)
            elif shp.shape_type == 13:  # 图片类型
                # 复制图片
                img = shp.image
                image_bytes = img.blob
                new_slide.shapes.add_picture(
                    image_bytes, shp.left, shp.top, shp.width, shp.height
                )

        for rpi, rpj in reversed_page_pptx_type.items():
            for shape in new_slide.shapes:
                shape.text_frame.clear()
                if shape.name == rpi:
                    new_text = rj[rpj]
                    for line in new_text.split('\n'):
                        p = shape.text_frame.add_paragraph()
                        p.text = line
    prs_target.save(output_filepath)


output_pptx(model_res1, pptx_demo)
# start()