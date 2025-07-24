from pptx import Presentation


from pptx import Presentation

def copy_pages(src_path, new_path):
    """
    按给定页码列表（从 1 开始）一页一页复制到新 PPT。
    页码可重复，用于多次插入同一页。
    """
    prs_src = Presentation(src_path)
    prs_new = Presentation(new_path)
    prs_new.slides[0] = prs_src.slides[0]
    # for page_idx, slide in enumerate(prs_new.slides):
    #     slide = prs_src.slides[0]
        # for shape_idx, shape in enumerate(slide.shapes):
            # print(shape_idx,shape)
            # print(shape.name)


    prs_new.save(r"C:\Users\Dell\Desktop\await_write1.pptx")



# 用法



if __name__ == "__main__":
    # 默认文件名
    source_pptx = "./files/templates/pptx/demo.pptx"
    # target_pptx = "./output/pptx/output.pptx"
    target_pptx = r"C:\Users\Dell\Desktop\await_write.pptx"
    # copy_slides(source_pptx, target_pptx)
    copy_pages(source_pptx,target_pptx)