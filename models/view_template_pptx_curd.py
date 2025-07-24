#!/usr/bin/env python
# -*- coding: UTF-8 -*-
'''
@Project ：AutoPPT
@File    ：view_template_pptx_curd.py
@Date    ：2025/7/23 15:15
@Descrip ：
'''


from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

import random
import json,os
from models.config_read import read_env,read_json_file

def analyze_pptx(pptx_path):
    """
    分析PPTX文件，提取每页的占位符信息
    返回格式: {页数: {"占位符内容": "占位符名"}}
    """
    prs = Presentation(pptx_path)
    result = {}
    for page_num, slide in enumerate(prs.slides, start=1):
        slide_data = {}
        # 遍历所有形状，提取占位符信息
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name:
                # 获取占位符名称和内容
                placeholder_name = shape.name
                placeholder_text = shape.text.strip()
                # 仅当内容不为空时记录
                if placeholder_text:
                    slide_data[placeholder_text] = placeholder_name
        # 将当前页的数据添加到结果中
        result[str(page_num)] = slide_data
    return result


## 找对应关系
"""
dict_public_placeholders = {"一级标题1": "fisrt_title1","一级标题2": "fisrt_title2"}
dict_detail = {'第 一 章 操作系统概论': '标题 1', '主讲老师：***': '文本框 11'}
dict_pptx_json_page = {'name': '章节主页', 'type': '章节主页', 'fisrt_title1': '标题 1'}
d_v = "标题 1"
"""
def get_placeholder_key(dict_public_placeholders, dict_pptx_json_page , d_v, converse=True):
    """
    通过多层字典映射查找最终的占位符键名

    参数:
    dict_detail: 存储详细内容的字典
    dict_pptx_json_page: 存储页面结构的字典
    dict_public_placeholders: 存储占位符映射的字典
    d_v: 需要查找的目标值（示例中为"标题 1"）

    返回:
    对应的占位符键名，如果未找到则返回 None
    """
    ## 显示
    if converse==True:
        # 步骤12: 使用找到的键作为dict_pptx_json_page的值，查找对应的键
        # 修正：这里应该直接使用d_v作为dict_pptx_json_page的值来查找键
        page_key = None
        for key, value in dict_pptx_json_page.items():
            if value == d_v:  # 直接使用d_v匹配dict_pptx_json_page中的值
                page_key = key
                break
        if page_key is None:
            return None  # 未找到匹配项
        # 步骤3: 使用找到的键作为dict_public_placeholders的值，查找对应的键
        placeholder_key = None
        for key, value in dict_public_placeholders.items():
            if value == page_key:
                placeholder_key = key
                break
        return placeholder_key
    ## 保存
    else:
        return None


def get_dict_pptx_page_params():
    return {
        "名称": "name",
        "类型": "type",
        "描述": "descript",
        "placeholders": {
            "一级标题1": "first_title1",
            "一级标题2": "first_title2",
            "二级标题1": "second_title1",
            "二级标题2": "second_title2",
            "三级标题1": "third_title1",
            "三级标题2": "third_title2",
            "正文内容1": "context1",
            "正文内容2": "context2",
            "插图1": "pic1",
            "插图2": "pic2"
        }
    }


def obj_searchall():
    dict_pptx_page_params = get_dict_pptx_page_params()
    ## 反向映射.{'fisrt_title1': '一级标题1', 'fisrt_title2': '一级标题2', 'second_title1': '二级标题1', 'second_title2': '二级标题2'......}
    dict_reversed_pptx_page_placeholders = {v: k for k, v in dict_pptx_page_params["placeholders"].items()}
    dict_env = read_env()
    dir_pptx = dict_env.get("dir_pptx")
    file_json_pptx = dict_env.get("file_json_pptx")
    global_demo_pptx_codenum = int(dict_env.get("global_demo_pptx_codenum"))
    dict_files_data = {}
    pptx_json = read_json_file(file_json_pptx)
    for filepath_name in os.listdir(dir_pptx):
        if filepath_name.endswith(".pptx"):
            filepath_pptx = os.path.join(dir_pptx, filepath_name)
            file_name = filepath_name.replace(".pptx","")
        else:
            continue
        ## 文件内容.'1': {'第 一 章 操作系统概论': '标题 1', '主讲老师：***': '文本框 11'}, '2': {'目录\nCONTENTS': '文本框 4', '第一节   操作系统的概念\n第二节   操作系统的发展\n第三节   操作系统分类\n第四节   操作系统设计\n第五节   操作系统启动': '文本框 11'}......}
        pptx_detail = analyze_pptx(filepath_pptx)
        file_pptx_json = pptx_json.get(file_name) if pptx_json.get(file_name) else {str(i): {} for i in range(100)}
        dict_file_data = {}
        for page_num, page_data in pptx_detail.items():
            # page_data: {'第 一 章 操作系统概论': '标题 1', '主讲老师：***': '文本框 11'}
            dict_file_pptx_json_page_num = file_pptx_json[page_num]
            # dict_reversed_file_pptx_json_page_num: {'章节主页': 'type', '': 'pic2', '标题 1': 'fisrt_title1', '文本框 11': 'context1'}
            dict_reversed_file_pptx_json_page_num = {v: k for k, v in dict_file_pptx_json_page_num.items()}
            dict_placeholders = {}
            for placeholder_text, placeholder_code in page_data.items():
                placeholder_type = dict_reversed_file_pptx_json_page_num.get(placeholder_code)
                placeholder_type_cn = dict_reversed_pptx_page_placeholders.get(placeholder_type)
                placeholder_text_deal = str(placeholder_text)[:global_demo_pptx_codenum]+"......" if len(placeholder_text)>global_demo_pptx_codenum else placeholder_text
                dict_placeholders[placeholder_text_deal] = {
                    "placeholder_code": placeholder_code,
                    "placeholder_type": placeholder_type,
                    "placeholder_type_cn": placeholder_type_cn,
                }
            dict_file_data[page_num] = {
                "name": dict_file_pptx_json_page_num["name"] if dict_file_pptx_json_page_num.get("name") else "",
                "type": dict_file_pptx_json_page_num["type"] if dict_file_pptx_json_page_num.get("type") else "",
                "descript": dict_file_pptx_json_page_num["descript"] if dict_file_pptx_json_page_num.get("descript") else "",
                "placeholders": dict_placeholders
            }
        dict_files_data[file_name] = dict_file_data
    return dict_files_data


def temp_obj_searchall():
    dict_env = read_env()
    dir_pptx = dict_env.get("dir_pptx")
    file_json_pptx = dict_env.get("file_json_pptx")
    dict_data = {}
    dict_data["file_name"] = []
    dict_data["detail"] = {}
    for filename in os.listdir(dir_pptx):
        if filename.endswith(".pptx"):
            filepath_pptx = os.path.join(dir_pptx, filename)
        else:
            continue
        dict_data["file_name"].append(filename.replace(".pptx",""))
        dict_data["pptx_json"] = read_json_file(file_json_pptx)
        dict_data["detail"][filename.replace(".pptx","")] = analyze_pptx(filepath_pptx)
    return dict_data


def obj_create(data):
    pass


def obj_update(data):
    dict_pptx_page_params = get_dict_pptx_page_params()
    dict_env = read_env()
    file_json_pptx = dict_env.get("file_json_pptx")
    pptx_json = read_json_file(file_json_pptx)
    ## {'demo': {'1': {'name': '章节主页', 'type': '章节主页', 'descript': '', 'first_title1': '标题 1', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '矩形 5', 'context2': '', 'pic1': '', 'pic2': ''}, '2': {'name': '目录页', 'type': '目录页', 'descript': '', 'first_title1': '', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 11', 'context2': '', 'pic1': '', 'pic2': ''}, '3': {'name': '学习要求', 'type': '学习要求', 'descript': '', 'first_title1': '', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 7', 'context2': '', 'pic1': '', 'pic2': ''}, '4': {'name': '重难点', 'type': '重难点', 'descript': '', 'first_title1': '', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '内容占位符 5', 'context2': '', 'pic1': '', 'pic2': ''}, '5': {'name': '节标题页', 'type': '节标题页', 'descript': '', 'first_title1': '', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 7', 'context2': '', 'pic1': '', 'pic2': ''}, '6': {'name': '普通正文', 'type': '普通正文', 'descript': '', 'first_title1': '', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 3', 'context2': '', 'pic1': '', 'pic2': ''}, '7': {'name': '小标题正文', 'type': '小标题正文', 'descript': '', 'first_title1': '标题 1', 'first_title2': '', 'second_title1': '文本框 10', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 4', 'context2': '', 'pic1': '', 'pic2': ''}, '8': {'name': '图片正文', 'type': '图片正文', 'descript': '', 'first_title1': '标题 1', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 9', 'context2': '', 'pic1': '图片占位符 20', 'pic2': ''}, '9': {'name': '小标题图片正文', 'type': '小标题图片正文', 'descript': '', 'first_title1': '标题 1', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 9', 'context2': '', 'pic1': '图片占位符 20', 'pic2': ''}, '10': {'name': '例题', 'type': '例题', 'descript': '', 'first_title1': '标题 1', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 10', 'context2': '', 'pic1': '', 'pic2': ''}, '11': {'name': '思维导图', 'type': '思维导图', 'descript': '', 'first_title1': '标题 1', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '图片占位符 20', 'context2': '', 'pic1': '', 'pic2': ''}, '12': {'name': '', 'main_title': '标题 1', 'small_title': '', 'context': '文本框 7', 'pic': '', 'type': '', 'descript': '', 'first_title1': '', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '', 'context2': '', 'pic1': '', 'pic2': ''}, '13': {'name': '尾页', 'type': '尾页', 'descript': '', 'first_title1': '标题 1', 'first_title2': '', 'second_title1': '', 'second_title2': '', 'third_title1': '', 'third_title2': '', 'context1': '文本框 7', 'context2': '', 'pic1': '', 'pic2': ''}}}
    file_name = data["file_name"]
    page_pptx_json = {}
    page_pptx_json["name"] = data["name"] if data.get("name") else ""
    page_pptx_json["type"] = data["type"] if data.get("type") else ""
    page_pptx_json["descript"] = data["descript"] if data.get("descript") else ""
    for data_k, data_v in data.items():
        if "placeholders" in data_k:
            placeholder_type = dict_pptx_page_params["placeholders"][data_v]
            page_pptx_json[placeholder_type] = data["placeholder_kv"].get(data_k)
    for i in list(dict_pptx_page_params['placeholders'].values()):
        if not page_pptx_json.get(i):
            page_pptx_json[i] = ""
    if not pptx_json.get(file_name):
        pptx_json[file_name] = {}
    pptx_json[file_name][str(data["page"])] = page_pptx_json
    with open(file_json_pptx, 'w', encoding='utf-8') as f:
        json.dump(pptx_json, f, ensure_ascii=False, indent=2)

def obj_delete(filename):
    dict_env = read_env()
    dir_pptx = dict_env.get("dir_pptx")
    file_json_pptx = dict_env.get("file_json_pptx")
    file_path = os.path.join(dir_pptx, filename+".pptx")
    os.remove(file_path) if os.path.isfile(file_path) else None
    pptx_json = read_json_file(file_json_pptx)
    pptx_json.pop(filename)
    with open(file_json_pptx, 'w', encoding='utf-8') as f:
        json.dump(pptx_json, f, ensure_ascii=False, indent=2)


